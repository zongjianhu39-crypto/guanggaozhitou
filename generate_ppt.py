#!/usr/bin/env python3
"""Generate the internal presentation for 交个朋友·广告智投工作台"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK    = RGBColor(0x0A, 0x0E, 0x27)
BG_SECTION = RGBColor(0x0F, 0x13, 0x2E)
ACCENT     = RGBColor(0x63, 0x66, 0xF1)
ACCENT2    = RGBColor(0x3B, 0x82, 0xF6)
ACCENT3    = RGBColor(0x22, 0xC5, 0x5E)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG    = RGBColor(0x1E, 0x29, 0x3B)
WARM       = RGBColor(0xF5, 0x9E, 0x0B)
RED_ACCENT = RGBColor(0xEF, 0x44, 0x44)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=BG_DARK):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_card(slide, left, top, width, height, title, body_lines,
             accent_color=ACCENT, title_size=16, body_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    shape.shadow.inherit = False

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent_color
    accent_bar.line.fill.background()

    add_text(slide, left + Inches(0.25), top + Inches(0.2),
             width - Inches(0.5), Inches(0.4),
             title, font_size=title_size, color=accent_color, bold=True)

    if body_lines:
        add_bullet_list(slide, left + Inches(0.25), top + Inches(0.65),
                        width - Inches(0.5), height - Inches(0.85),
                        body_lines, font_size=body_size, color=LIGHT_GRAY, spacing=Pt(4))


def add_number_card(slide, left, top, width, height, number, label, accent=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    add_text(slide, left, top + Inches(0.15), width, Inches(0.5),
             number, font_size=36, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.7), width, Inches(0.3),
             label, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def add_divider(slide, top, color=ACCENT, left=Inches(0.8), width=None):
    if width is None:
        width = W - Inches(1.6)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


# ═══════════════════════════════════════════
# SLIDE 1 – Cover
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-1.5), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x14, 0x1A, 0x3A)
circle.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(4), Inches(4), Inches(4))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x10, 0x15, 0x30)
circle2.line.fill.background()

add_text(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(1),
         "交个朋友 · 广告智投工作台", font_size=44, bold=True, align=PP_ALIGN.CENTER)

add_divider(slide, Inches(3.0), ACCENT, Inches(4.5), Inches(4.3))

add_text(slide, Inches(1.2), Inches(3.3), Inches(11), Inches(0.6),
         "用数据驱动投放决策，用 AI 沉淀运营经验", font_size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.2), Inches(4.3), Inches(11), Inches(0.5),
         "内部产品介绍  ·  2026 年 4 月", font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

cards_data = [
    ("Monitor", "实时盯盘", ACCENT),
    ("Explain", "AI 解读", ACCENT2),
    ("Archive", "经验归档", ACCENT3),
    ("Tune", "策略优化", WARM),
]
card_w = Inches(2.2)
card_h = Inches(0.7)
total_w = card_w * 4 + Inches(0.3) * 3
start_x = (W - total_w) // 2
for i, (en, zh, clr) in enumerate(cards_data):
    x = start_x + i * (card_w + Inches(0.3))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.3), card_w, card_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = clr
    shape.line.width = Pt(1.5)
    add_text(slide, x, Inches(5.35), card_w, Inches(0.35),
             en, font_size=14, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.65), card_w, Inches(0.3),
             zh, font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 2 – Pain Points / Why
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_SECTION)

add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6),
         "为什么需要这个工作台？", font_size=32, bold=True)
add_divider(slide, Inches(1.15), ACCENT, Inches(0.8), Inches(3))

add_text(slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4),
         "投放团队每天面临的核心挑战", font_size=16, color=LIGHT_GRAY)

pain_points = [
    ("数据散落", "投放数据分布在万相台、代理后台、短直联动、财务表等多个平台，每次汇总需要反复切换", RED_ACCENT),
    ("人工盯盘低效", "投手每天花大量时间在 Excel 里比对花费、ROI、成本，容易遗漏异常", WARM),
    ("经验难沉淀", "好的投放策略和踩过的坑停留在个人经验层面，新人上手慢", ACCENT2),
    ("口径不统一", "同一个指标不同人有不同理解，对账时经常出现偏差", ACCENT),
    ("计划管理粗放", "月度计划拆解靠 Excel 手传，实际与计划的偏差发现滞后", ACCENT3),
]

for i, (title, desc, clr) in enumerate(pain_points):
    y = Inches(2.0) + i * Inches(1.05)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.08), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = clr
    dot.line.fill.background()
    add_text(slide, Inches(1.35), y - Inches(0.05), Inches(4.5), Inches(0.35),
             title, font_size=18, color=clr, bold=True)
    add_text(slide, Inches(1.35), y + Inches(0.3), Inches(4.8), Inches(0.5),
             desc, font_size=13, color=LIGHT_GRAY)

sol_x = Inches(7.2)
sol_w = Inches(5.5)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sol_x, Inches(1.6), sol_w, Inches(5.2))
shape.fill.solid()
shape.fill.fore_color.rgb = CARD_BG
shape.line.fill.background()

add_text(slide, sol_x + Inches(0.3), Inches(1.8), sol_w - Inches(0.6), Inches(0.5),
         "我们的解法", font_size=22, color=ACCENT3, bold=True)

solutions = [
    "一个平台聚合所有投放数据源",
    "自动计算 KPI，异常实时可见",
    "AI 自动生成分析报告 + 经验库",
    "统一指标口径，公开透明",
    "计划拆解在线协同，实时对账",
    "自然语言问数，不用写 SQL",
]
for i, sol in enumerate(solutions):
    y = Inches(2.5) + i * Inches(0.65)
    check = slide.shapes.add_shape(MSO_SHAPE.OVAL, sol_x + Inches(0.4), y + Inches(0.05), Inches(0.2), Inches(0.2))
    check.fill.solid()
    check.fill.fore_color.rgb = ACCENT3
    check.line.fill.background()
    add_text(slide, sol_x + Inches(0.4), y + Inches(0.02), Inches(0.2), Inches(0.2),
             "✓", font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, sol_x + Inches(0.8), y, sol_w - Inches(1.2), Inches(0.35),
             sol, font_size=15, color=WHITE)


# ═══════════════════════════════════════════
# SLIDE 3 – Overview / What Is It
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
         "工作台全景", font_size=32, bold=True)
add_divider(slide, Inches(1.15), ACCENT, Inches(0.8), Inches(2.5))

add_text(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
         "7 大核心模块，覆盖投放全链路", font_size=18, color=LIGHT_GRAY)

modules = [
    ("数据看板", "投放分析 · 人群 · 单品\n多维度 KPI 一屏汇总", ACCENT, "核心盯盘"),
    ("AI 智能分析", "一键生成日报/周报\n自动识别异常与风险", ACCENT2, "AI 驱动"),
    ("洞察中心", "报告归档 · 风险追踪\n经验沉淀与复盘", ACCENT3, "知识管理"),
    ("GenBI 智能问数", "自然语言提问\n受控查询，不编造数据", WARM, "自助分析"),
    ("计划拆解", "月度计划在线管理\n实际 vs 计划实时对比", RGBColor(0xA7, 0x8B, 0xFA), "协同管理"),
    ("指标规则台", "统一指标口径\n公式与规则透明公开", RGBColor(0x06, 0xB6, 0xD4), "口径统一"),
    ("Prompt 管理", "AI 行为精细调控\n六大类 Prompt 版本管理", RGBColor(0xEC, 0x48, 0x99), "AI 治理"),
]

col_count = 4
row1 = modules[:4]
row2 = modules[4:]

for row_idx, row_modules in enumerate([row1, row2]):
    count = len(row_modules)
    cw = Inches(2.7)
    ch = Inches(2.2)
    gap = Inches(0.3)
    total = cw * count + gap * (count - 1)
    sx = (W - total) // 2
    sy = Inches(2.1) + row_idx * (ch + Inches(0.35))

    for i, (name, desc, clr, tag) in enumerate(row_modules):
        x = sx + i * (cw + gap)
        add_card(slide, x, sy, cw, ch, name,
                 desc.split('\n'), accent_color=clr, title_size=17, body_size=13)
        tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           x + cw - Inches(1.1), sy + ch - Inches(0.4),
                                           Inches(0.9), Inches(0.25))
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = clr
        tag_shape.line.fill.background()
        add_text(slide, x + cw - Inches(1.1), sy + ch - Inches(0.4),
                 Inches(0.9), Inches(0.25),
                 tag, font_size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 4 – Data Dashboard Deep Dive
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_SECTION)

add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.6),
         "数据看板", font_size=32, bold=True, color=ACCENT)
add_text(slide, Inches(0.8), Inches(1.1), Inches(6), Inches(0.4),
         "投放数据的「中央驾驶舱」", font_size=18, color=LIGHT_GRAY)
add_divider(slide, Inches(1.55), ACCENT, Inches(0.8), Inches(2.5))

tabs = [
    ("投放分析", "花费、ROI、盈亏平衡 ROI、去退 ROI\n支持分月/分周/每日多粒度汇总\n自动高亮异常指标", ACCENT),
    ("人群维度", "按人群分层查看投放效果\n精准定位高消耗低 ROI 人群\n辅助人群策略调整", ACCENT2),
    ("单品广告", "按商品维度拆解投放效果\n花费 & ROI 一目了然\n快速识别弱品 / 明星品", ACCENT3),
]

for i, (name, desc, clr) in enumerate(tabs):
    x = Inches(0.8) + i * Inches(4.1)
    add_card(slide, x, Inches(1.9), Inches(3.8), Inches(2.5),
             name, desc.split('\n'), accent_color=clr, title_size=18, body_size=14)

features_left = [
    "日期范围灵活选择，支持跨月对比",
    "KPI 卡片快速定位关键指标",
    "CSV 一键导出（日报 / 全量 / 单品）",
]
features_right = [
    "数据源自动聚合，无需手动拼表",
    "骨架屏加载，体验流畅",
    "飞书登录集成，安全便捷",
]

add_text(slide, Inches(0.8), Inches(4.7), Inches(4), Inches(0.3),
         "核心能力", font_size=16, color=WHITE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(5.1), Inches(5.5), Inches(2),
                features_left, font_size=14, color=LIGHT_GRAY)
add_bullet_list(slide, Inches(6.8), Inches(5.1), Inches(5.5), Inches(2),
                features_right, font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════
# SLIDE 5 – AI Analysis + Insights Center
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
         "AI 智能分析 + 洞察中心", font_size=32, bold=True, color=ACCENT2)
add_divider(slide, Inches(1.15), ACCENT2, Inches(0.8), Inches(3))

# Left: AI Analysis
add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3),
         "AI 智能分析", [], accent_color=ACCENT2, title_size=20)

ai_steps = [
    ("1", "点击「AI 分析」", "在数据看板中选择日期和维度"),
    ("2", "自动拉取数据", "聚合投放、财务、人群等多源数据"),
    ("3", "大模型生成报告", "基于 Prompt 模板 + 业务红线分析"),
    ("4", "一键发布归档", "报告自动落库到洞察中心"),
]

for i, (num, title, desc) in enumerate(ai_steps):
    y = Inches(2.2) + i * Inches(1.0)
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.35), Inches(0.35))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = ACCENT2
    num_shape.line.fill.background()
    add_text(slide, Inches(1.2), y + Inches(0.02), Inches(0.35), Inches(0.35),
             num, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.75), y - Inches(0.02), Inches(4), Inches(0.3),
             title, font_size=16, color=WHITE, bold=True)
    add_text(slide, Inches(1.75), y + Inches(0.3), Inches(4.5), Inches(0.3),
             desc, font_size=12, color=LIGHT_GRAY)

# Right: Insights Center
add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.3),
         "洞察中心", [], accent_color=ACCENT3, title_size=20)

insight_features = [
    ("报告库", "所有 AI 报告按日期、类型归档\n支持风险等级筛选"),
    ("结构化详情", "焦点问题 · 高消耗人群 · 执行动作\n财务修正 · 重点场次一网打尽"),
    ("经验沉淀", "从历史报告中提炼经验\n新人可以快速学习投放策略"),
    ("知识复用", "经验库与 GenBI 联动\n每次分析自动引用历史经验"),
]

for i, (title, desc) in enumerate(insight_features):
    y = Inches(2.2) + i * Inches(1.15)
    add_text(slide, Inches(7.4), y, Inches(5), Inches(0.3),
             f"▸ {title}", font_size=16, color=ACCENT3, bold=True)
    for j, line in enumerate(desc.split('\n')):
        add_text(slide, Inches(7.6), y + Inches(0.32) + j * Inches(0.28),
                 Inches(4.8), Inches(0.3),
                 line, font_size=12, color=LIGHT_GRAY)


# ═══════════════════════════════════════════
# SLIDE 6 – GenBI
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_SECTION)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
         "GenBI · 智能问数", font_size=32, bold=True, color=WARM)
add_text(slide, Inches(0.8), Inches(1.1), Inches(8), Inches(0.4),
         "用自然语言提问，不用写 SQL，也不用翻多个后台", font_size=18, color=LIGHT_GRAY)
add_divider(slide, Inches(1.55), WARM, Inches(0.8), Inches(3))

add_card(slide, Inches(0.8), Inches(1.9), Inches(6), Inches(1.6),
         "你可以这样问", [], accent_color=WARM, title_size=18)

questions = [
    "「昨天各人群的花费占比是多少？」",
    "「上周哪些商品 ROI 低于盈亏平衡线？」",
    "「帮我看看这个月老客和新客的成交对比」",
    "「上周花费比上上周波动了多少？」",
]
add_bullet_list(slide, Inches(1.2), Inches(2.55), Inches(5.2), Inches(1.5),
                questions, font_size=14, color=LIGHT_GRAY, spacing=Pt(4))

# Right side: How it works
add_card(slide, Inches(7.2), Inches(1.9), Inches(5.5), Inches(4.8),
         "工作原理", [], accent_color=WARM, title_size=18)

how_items = [
    ("意图识别", "NLP 解析问题 → 匹配预定义意图\n（人群预算、弱品分析、花费波动等）"),
    ("受控查询", "只查真实数仓数据\n不允许 AI 编造数字，保证准确性"),
    ("智能引用", "自动关联指标规则、历史报告\n返回带参考来源的结构化结果"),
    ("结果呈现", "表格 + 结论文本 + 参考链接\n高亮关键标签，直观易懂"),
]

for i, (title, desc) in enumerate(how_items):
    y = Inches(2.6) + i * Inches(1.1)
    add_text(slide, Inches(7.6), y, Inches(4.8), Inches(0.3),
             title, font_size=16, color=WARM, bold=True)
    for j, line in enumerate(desc.split('\n')):
        add_text(slide, Inches(7.6), y + Inches(0.3) + j * Inches(0.26),
                 Inches(4.8), Inches(0.3),
                 line, font_size=12, color=LIGHT_GRAY)

add_card(slide, Inches(0.8), Inches(3.8), Inches(6), Inches(2.9),
         "核心优势", [], accent_color=WARM, title_size=18)

advantages = [
    "不需要任何技术背景，自然语言即可查数",
    "受控意图机制，杜绝 AI「幻觉」编造数据",
    "结果自带参考来源，可溯源验证",
    "覆盖 10+ 常见分析意图，持续扩展",
    "与指标规则台对齐，口径一致",
]
add_bullet_list(slide, Inches(1.2), Inches(4.5), Inches(5.2), Inches(2.5),
                advantages, font_size=14, color=LIGHT_GRAY, spacing=Pt(5))


# ═══════════════════════════════════════════
# SLIDE 7 – Plan Dashboard
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
         "计划拆解", font_size=32, bold=True, color=RGBColor(0xA7, 0x8B, 0xFA))
add_text(slide, Inches(0.8), Inches(1.1), Inches(8), Inches(0.4),
         "让月度投放计划从 Excel 走向在线协同", font_size=18, color=LIGHT_GRAY)
add_divider(slide, Inches(1.55), RGBColor(0xA7, 0x8B, 0xFA), Inches(0.8), Inches(3))

plan_features = [
    ("按月管理", "月份快速切换，一键查看整月计划\n万相台 + 有客代投双通道计划金额", RGBColor(0xA7, 0x8B, 0xFA)),
    ("实时对账", "实际花费自动聚合三大数据源\n万相台 + 有客代投 + 短直联动", ACCENT2),
    ("活动时间轴", "大促 / 日常活动可视化\n活动节奏与计划金额联动", ACCENT3),
    ("同比参考", "2025 年同期花费、成交、预售、毛利\n历史数据辅助决策", WARM),
]

for i, (title, desc, clr) in enumerate(plan_features):
    x = Inches(0.8) + (i % 2) * Inches(6.2)
    y = Inches(1.9) + (i // 2) * Inches(2.5)
    add_card(slide, x, y, Inches(5.8), Inches(2.2),
             title, desc.split('\n'), accent_color=clr, title_size=18, body_size=14)

more = [
    "行内编辑 + 批量保存草稿",
    "备注悬浮预览 + 点击编辑",
    "CSV 一键导出",
    "合计栏 + 完成率实时统计",
]
add_text(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
         "  ·  ".join(more), font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 8 – Prompt Management + Metric Rules
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_SECTION)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
         "AI 治理 · 指标口径", font_size=32, bold=True, color=RGBColor(0xEC, 0x48, 0x99))
add_divider(slide, Inches(1.15), RGBColor(0xEC, 0x48, 0x99), Inches(0.8), Inches(3))

# Left: Prompt Management
add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.3),
         "Prompt 管理", [], accent_color=RGBColor(0xEC, 0x48, 0x99), title_size=20)

prompt_cats = [
    ("长期记忆", "定义 AI 的身份认知和基本原则"),
    ("灵魂 Prompt", "核心行为规范与风格调性"),
    ("技能模块", "特定场景的分析能力配置"),
    ("数据分析", "数据解读的方法论和框架"),
    ("运营业务", "业务逻辑和运营规则设定"),
    ("业务红线", "不可触碰的底线规则"),
]

for i, (cat, desc) in enumerate(prompt_cats):
    y = Inches(2.2) + i * Inches(0.72)
    tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(1.2), y, Inches(1.2), Inches(0.3))
    tag_shape.fill.solid()
    tag_shape.fill.fore_color.rgb = RGBColor(0xEC, 0x48, 0x99)
    tag_shape.line.fill.background()
    add_text(slide, Inches(1.2), y, Inches(1.2), Inches(0.3),
             cat, font_size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(2.6), y, Inches(3.5), Inches(0.3),
             desc, font_size=13, color=LIGHT_GRAY)

add_text(slide, Inches(1.2), Inches(6.0), Inches(5), Inches(0.4),
         "支持草稿 → 发布 → 回滚，版本可控", font_size=13, color=RGBColor(0xEC, 0x48, 0x99))

# Right: Metric Rules
add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.3),
         "指标与规则台", [], accent_color=RGBColor(0x06, 0xB6, 0xD4), title_size=20)

metric_items = [
    "所有看板指标的计算公式公开可查",
    "人群规则定义清晰可溯",
    "GenBI 语义配置透明展示",
    "带搜索功能，快速定位指标",
    "数据看板规则 / GenBI 语义双 Tab",
    "确保团队内口径一致，减少对账偏差",
]
for i, item in enumerate(metric_items):
    y = Inches(2.2) + i * Inches(0.72)
    check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.4), y + Inches(0.05), Inches(0.2), Inches(0.2))
    check.fill.solid()
    check.fill.fore_color.rgb = RGBColor(0x06, 0xB6, 0xD4)
    check.line.fill.background()
    add_text(slide, Inches(7.4), y + Inches(0.02), Inches(0.2), Inches(0.2),
             "✓", font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(7.8), y, Inches(4.5), Inches(0.35),
             item, font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════
# SLIDE 9 – Architecture
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
         "技术架构", font_size=32, bold=True)
add_divider(slide, Inches(1.15), ACCENT, Inches(0.8), Inches(2.5))

layers = [
    ("用户层", "飞书 OAuth 登录 · 响应式 Web · 移动端适配", ACCENT,
     ["投手", "运营", "管理者"]),
    ("功能层", "数据看板 · AI 分析 · GenBI · 计划拆解 · 洞察中心 · Prompt 管理 · 指标规则", ACCENT2,
     ["7 大功能模块"]),
    ("服务层", "Supabase Edge Functions (Deno) · 统一鉴权 · MiniMax AI · Prompt 模板引擎", ACCENT3,
     ["8 个 Edge Functions"]),
    ("数据层", "万相台 · 代理数据 · 短直联动 · 财务数据 · 淘宝直播 · 单品广告 · AI 报告库", WARM,
     ["10+ 数据表"]),
]

for i, (name, desc, clr, tags) in enumerate(layers):
    y = Inches(1.6) + i * Inches(1.4)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.8), y, Inches(11.7), Inches(1.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = clr
    shape.line.width = Pt(1)

    add_text(slide, Inches(1.1), y + Inches(0.1), Inches(2), Inches(0.35),
             name, font_size=18, color=clr, bold=True)
    add_text(slide, Inches(1.1), y + Inches(0.5), Inches(8), Inches(0.35),
             desc, font_size=12, color=LIGHT_GRAY)

    for j, tag in enumerate(tags):
        tx = Inches(9.5) + j * Inches(1.1)
        tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           tx, y + Inches(0.35), Inches(1.5), Inches(0.35))
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = clr
        tag_shape.line.fill.background()
        add_text(slide, tx, y + Inches(0.35), Inches(1.5), Inches(0.35),
                 tag, font_size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    if i < len(layers) - 1:
        arrow_y = y + Inches(1.15)
        arrow = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                       Inches(6.4), arrow_y + Inches(0.02),
                                       Inches(0.3), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = clr
        arrow.line.fill.background()
        arrow.rotation = 180.0


# ═══════════════════════════════════════════
# SLIDE 10 – Key Numbers
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_SECTION)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
         "关键数字", font_size=32, bold=True)
add_divider(slide, Inches(1.15), ACCENT, Inches(0.8), Inches(2.5))

numbers = [
    ("7", "功能模块", ACCENT),
    ("8", "Edge Functions", ACCENT2),
    ("10+", "数据表", ACCENT3),
    ("3", "用户角色", WARM),
    ("10+", "GenBI 意图", RGBColor(0xA7, 0x8B, 0xFA)),
    ("6", "Prompt 分类", RGBColor(0xEC, 0x48, 0x99)),
]

cw = Inches(1.8)
ch = Inches(1.2)
gap = Inches(0.3)
total = cw * 6 + gap * 5
sx = (W - total) // 2
for i, (num, label, clr) in enumerate(numbers):
    add_number_card(slide, sx + i * (cw + gap), Inches(1.8), cw, ch, num, label, clr)

benefits = [
    ("告别多平台切换", "所有投放数据统一入口，节省投手每日 30~60 分钟汇总时间"),
    ("AI 赋能日常复盘", "一键生成分析报告，自动归档洞察中心，沉淀团队经验"),
    ("计划执行可视化", "从 Excel 迁移到在线协同，计划偏差即时可见，提升执行力"),
    ("口径统一透明", "指标公式公开可查，消除团队内对账偏差，降低沟通成本"),
]

for i, (title, desc) in enumerate(benefits):
    x = Inches(0.8) + (i % 2) * Inches(6.2)
    y = Inches(3.5) + (i // 2) * Inches(1.7)
    add_card(slide, x, y, Inches(5.8), Inches(1.5),
             title, [desc], accent_color=[ACCENT, ACCENT2, ACCENT3, WARM][i],
             title_size=18, body_size=14)


# ═══════════════════════════════════════════
# SLIDE 11 – Role-Based Value
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
         "不同角色的使用价值", font_size=32, bold=True)
add_divider(slide, Inches(1.15), ACCENT, Inches(0.8), Inches(2.5))

roles = [
    ("投手", "日常高频使用者",
     ["数据看板盯盘 → 发现异常立即调整",
      "GenBI 快速查数 → 不用等 BI 同事排期",
      "计划拆解对账 → 实时掌握预算执行进度",
      "AI 日报 → 下班前一键生成当日复盘"],
     ACCENT),
    ("运营",  "策略制定与优化者",
     ["洞察中心复盘 → 从报告中提炼运营策略",
      "经验沉淀 → 好的策略归档为团队资产",
      "人群 & 单品维度 → 辅助选品和人群策略",
      "指标规则 → 统一团队分析口径"],
     ACCENT2),
    ("管理者",  "全局把控与决策者",
     ["数据看板总览 → 一屏掌握整体投放健康度",
      "计划拆解 → 月度目标追踪，偏差一目了然",
      "AI 分析报告 → 快速获取结构化经营洞察",
      "Prompt 管理 → 确保 AI 输出符合业务规范"],
     ACCENT3),
]

cw = Inches(3.8)
ch = Inches(5.0)
gap = Inches(0.3)
total = cw * 3 + gap * 2
sx = (W - total) // 2

for i, (role, subtitle, items, clr) in enumerate(roles):
    x = sx + i * (cw + gap)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), cw, ch)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = clr
    shape.line.width = Pt(1.5)

    icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        x + (cw - Inches(0.6)) // 2, Inches(1.85),
                                        Inches(0.6), Inches(0.6))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = clr
    icon_shape.line.fill.background()
    icons = ["T", "O", "M"]
    add_text(slide, x + (cw - Inches(0.6)) // 2, Inches(1.92),
             Inches(0.6), Inches(0.5),
             icons[i], font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x, Inches(2.6), cw, Inches(0.4),
             role, font_size=22, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.0), cw, Inches(0.3),
             subtitle, font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        y = Inches(3.5) + j * Inches(0.72)
        parts = item.split(' → ')
        add_text(slide, x + Inches(0.3), y, cw - Inches(0.6), Inches(0.3),
                 f"▸ {parts[0]}", font_size=12, color=WHITE, bold=True)
        if len(parts) > 1:
            add_text(slide, x + Inches(0.5), y + Inches(0.28), cw - Inches(0.8), Inches(0.3),
                     parts[1], font_size=11, color=LIGHT_GRAY)


# ═══════════════════════════════════════════
# SLIDE 12 – Getting Started + CTA
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(3), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x14, 0x1A, 0x3A)
circle.line.fill.background()

add_text(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(0.8),
         "开始使用", font_size=36, bold=True, align=PP_ALIGN.CENTER)

add_divider(slide, Inches(2.4), ACCENT, Inches(5), Inches(3.3))

steps = [
    ("Step 1", "飞书扫码登录", "使用公司飞书账号即可登录，无需额外注册", ACCENT),
    ("Step 2", "进入数据看板", "选择日期范围，查看投放数据全景", ACCENT2),
    ("Step 3", "尝试 AI 分析", "点击 AI 分析按钮，体验一键报告生成", ACCENT3),
    ("Step 4", "探索更多功能", "GenBI 问数、计划拆解、洞察中心等你发现", WARM),
]

cw = Inches(2.7)
ch = Inches(1.8)
gap = Inches(0.3)
total = cw * 4 + gap * 3
sx = (W - total) // 2

for i, (step, title, desc, clr) in enumerate(steps):
    x = sx + i * (cw + gap)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.0), cw, ch)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = clr
    shape.line.width = Pt(1.5)

    add_text(slide, x, Inches(3.1), cw, Inches(0.3),
             step, font_size=12, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.4), cw, Inches(0.35),
             title, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.15), Inches(3.85), cw - Inches(0.3), Inches(0.6),
             desc, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(0.6),
         "访问地址：friends.wang", font_size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1.2), Inches(6.0), Inches(11), Inches(0.5),
         "有任何问题或建议，欢迎随时反馈！", font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── Save ──
output_path = "/Users/zhouhao/Desktop/website/广告智投工作台介绍.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
